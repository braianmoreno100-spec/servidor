"""
reporte_pptx.py
Genera el reporte mensual PowerPoint al estilo Inverfarma,
usando el template original (fondo, logo, diseño).
"""

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
from datetime import date
import calendar
from io import BytesIO
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from ..database import get_db
from ..models import Turno, Orden, Parada, Desperdicio, RegistroProduccion

router = APIRouter(tags=["reporte-pptx"])

# ── Rutas assets ──────────────────────────────────────────────────────────────
BASE_DIR      = os.path.dirname(os.path.abspath(__file__))
ASSETS_DIR    = os.path.normpath(os.path.join(BASE_DIR, "..", "assets"))
TEMPLATE_PATH = os.path.join(ASSETS_DIR, "template_inverfarma.pptx")
LOGO_PATH     = os.path.join(ASSETS_DIR, "logo_inverfarma.png")

MESES_ES = [
    '', 'ENERO','FEBRERO','MARZO','ABRIL','MAYO','JUNIO',
    'JULIO','AGOSTO','SEPTIEMBRE','OCTUBRE','NOVIEMBRE','DICIEMBRE'
]

# ── Helpers OEE ───────────────────────────────────────────────────────────────
def _parse_hora(h: str) -> float:
    if not h:
        return 0.0
    import re
    clean = h.replace(".", "").replace("  ", " ").strip().lower()
    m = re.match(r"(\d{1,2}):(\d{2})(?:\s*(am|pm))?", clean)
    if not m:
        return 0.0
    hh, mm, ampm = int(m.group(1)), int(m.group(2)), m.group(3)
    if ampm == "pm" and hh != 12:
        hh += 12
    if ampm == "am" and hh == 12:
        hh = 0
    return hh + mm / 60


def _calcular_oee(turno, orden) -> dict:
    contador     = sum(r.cantidad for r in turno.registros_produccion)
    min_np       = sum(p.minutos for p in turno.paradas if not p.programada)
    n_paradas_np = len([p for p in turno.paradas if not p.programada])

    inicio = _parse_hora(turno.hora_inicio)
    if turno.hora_fin:
        fin = _parse_hora(turno.hora_fin)
    else:
        from datetime import datetime
        now = datetime.now()
        fin = now.hour + now.minute / 60
    tiempo_real = fin - inicio
    if tiempo_real < 0:
        tiempo_real += 24
    tiempo_real = min(tiempo_real, 12)

    tiempo_prog = 12.0
    ciclo       = float(orden.ciclos or 0)
    cavidades   = int(orden.cavidades or 1)
    tiempo_disp = max(tiempo_real - min_np / 60, 0)
    disp        = (tiempo_disp / tiempo_prog * 100) if tiempo_prog > 0 else 0
    prod_plan   = (tiempo_prog * 3600 / ciclo * cavidades) if ciclo > 0 else 0
    rend        = (contador / prod_plan * 100) if prod_plan > 0 else 0
    oee         = (disp / 100) * (rend / 100) * 100

    return {
        "oee":           round(min(oee, 150), 2),
        "disponibilidad":round(min(disp, 150), 2),
        "rendimiento":   round(min(rend, 150), 2),
        "calidad":       100.0,
        "min_np":        min_np,
        "n_paradas_np":  n_paradas_np,
        "tiempo_real_h": round(tiempo_real, 2),
        "contador":      contador,
    }


# ── Recolección de datos ──────────────────────────────────────────────────────
def _datos_mes(mes: int, anio: int, db: Session) -> dict:
    primer_dia = str(date(anio, mes, 1))
    ultimo_dia = str(date(anio, mes, calendar.monthrange(anio, mes)[1]))

    turnos = db.query(Turno).filter(
        Turno.fecha >= primer_dia,
        Turno.fecha <= ultimo_dia,
    ).all()

    orden_ids   = list({t.orden_id for t in turnos})
    ordenes     = db.query(Orden).filter(Orden.id.in_(orden_ids)).all()
    ordenes_map = {o.id: o for o in ordenes}

    por_tipo_mes: dict = {}
    desperdicios_mes: list = []
    paradas_data: list = []
    kg_mes: dict = {}

    for turno in turnos:
        orden = ordenes_map.get(turno.orden_id)
        if not orden:
            continue

        kpis = _calcular_oee(turno, orden)
        tipo = orden.tipo_maquina
        num  = str(orden.numero_maquina)
        oee  = kpis["oee"]

        if tipo == "linea" and num == "1":
            tipo_key = "linea_copro"
        elif tipo == "linea" and num == "2":
            tipo_key = "linea_orina"
        else:
            tipo_key = tipo

        try:
            mes_turno = int(turno.fecha[5:7])
        except:
            mes_turno = mes

        if tipo_key not in por_tipo_mes:
            por_tipo_mes[tipo_key] = {}
        if mes_turno not in por_tipo_mes[tipo_key]:
            por_tipo_mes[tipo_key][mes_turno] = []
        if oee > 0:
            por_tipo_mes[tipo_key][mes_turno].append(oee)

        for d in turno.desperdicios:
            desperdicios_mes.append((d.defecto, d.cantidad))

        paradas_data.append({
            "tipo_key":  tipo_key,
            "min_np":    kpis["min_np"],
            "n_paradas": kpis["n_paradas_np"],
            "tiempo_h":  kpis["tiempo_real_h"],
        })

        from ..models import Producto
        producto = db.query(Producto).filter(
            Producto.codigo == orden.codigo_producto
        ).first()
        if producto and producto.peso_pieza and kpis["contador"] > 0:
            cav = orden.cavidades or 1
            kg  = (kpis["contador"] * cav * producto.peso_pieza) / 1000
            maq_key = f"{tipo}_{num}"
            kg_mes[maq_key] = round(kg_mes.get(maq_key, 0) + kg, 2)

    mttr_acc: dict = {}
    for row in paradas_data:
        tk = row["tipo_key"]
        if tk not in mttr_acc:
            mttr_acc[tk] = {"min_np": 0, "n": 0, "h": 0}
        mttr_acc[tk]["min_np"] += row["min_np"]
        mttr_acc[tk]["n"]      += row["n_paradas"]
        mttr_acc[tk]["h"]      += row["tiempo_h"]

    mttr_result = {}
    for tk, d in mttr_acc.items():
        n    = d["n"]
        mttr = (d["min_np"] / n) if n > 0 else 0
        top  = d["h"] * 60 - d["min_np"]
        mtbf = (top / n) if n > 0 else d["h"] * 60
        disp = (top / (d["h"] * 60) * 100) if d["h"] > 0 else 0
        mttr_result[tk] = {
            "mttr": round(mttr, 1),
            "mtbf": round(mtbf, 1),
            "disp": round(disp, 1),
        }

    desp_agrup: dict = {}
    for defecto, cantidad in desperdicios_mes:
        desp_agrup[defecto] = desp_agrup.get(defecto, 0) + cantidad

    return {
        "por_tipo_mes": por_tipo_mes,
        "mttr_mtbf":    mttr_result,
        "desperdicios": desp_agrup,
        "kg_mes":       kg_mes,
    }


# ── Gráficas matplotlib ───────────────────────────────────────────────────────
def _setup_ax(ax, fig):
    fig.patch.set_facecolor('#EBEBEB')
    ax.set_facecolor('#DCDCDC')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#aaa')
    ax.spines['bottom'].set_color('#aaa')
    ax.grid(True, axis='y', color='white', linewidth=0.9, alpha=0.9)


def _chart_oee_linea(titulo: str, datos_mes: dict, tipo_key: str, mes_actual: int) -> BytesIO:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt

    meses_corto = ['ENE','FEB','MAR','ABR','MAY','JUN','JUL','AGO','SEP','OCT','NOV','DIC']
    fig, ax = plt.subplots(figsize=(9.4, 5.0), dpi=120)
    _setup_ax(ax, fig)

    y = []
    for m in range(1, 13):
        vals = datos_mes.get(tipo_key, {}).get(m, [])
        avg  = round(sum(vals)/len(vals), 2) if vals else None
        y.append(avg if m <= mes_actual else None)

    x_plot = [i for i, v in enumerate(y) if v is not None]
    y_plot = [v for v in y if v is not None]

    if x_plot:
        ax.plot(x_plot, y_plot, 's-', color='#70AD47', linewidth=2.5,
                markersize=7, label='2026', zorder=3)
        for xi, yi in zip(x_plot, y_plot):
            ax.annotate(f'{yi:.2f}', (xi, yi),
                        textcoords='offset points', xytext=(0, 9),
                        ha='center', fontsize=8.5, fontweight='bold', color='#2e6b0e')
        prom = sum(y_plot)/len(y_plot)
        ax.axhline(prom, color='#4472C4', linewidth=1.2, linestyle='--', alpha=0.7)
        ax.text(11.5, prom+0.5, f'Prom\n{prom:.1f}%', fontsize=7.5, color='#4472C4', ha='right')

    ax.set_xticks(range(12))
    ax.set_xticklabels(meses_corto, fontsize=9)
    if y_plot:
        margin = max(8, (max(y_plot)-min(y_plot))*0.3)
        ax.set_ylim(max(0, min(y_plot)-margin), min(110, max(y_plot)+margin))
    else:
        ax.set_ylim(0, 110)
    ax.set_title(titulo, fontsize=13, fontweight='bold', pad=10)
    if x_plot:
        ax.legend(loc='lower right', fontsize=9, framealpha=0.8)

    plt.tight_layout(pad=0.8)
    buf = BytesIO()
    plt.savefig(buf, format='png', dpi=120, bbox_inches='tight', facecolor='#EBEBEB')
    buf.seek(0)
    plt.close()
    return buf


def _chart_kg(kg_mes: dict, mes_nombre: str, anio: int) -> BytesIO:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt

    fig, ax = plt.subplots(figsize=(9.4, 4.8), dpi=120)
    _setup_ax(ax, fig)

    if not kg_mes:
        ax.text(0.5, 0.5, 'Sin datos de peso pieza configurado',
                ha='center', va='center', fontsize=12, color='#888', transform=ax.transAxes)
        ax.axis('off')
    else:
        color_map = {"inyeccion":"#4472C4","soplado":"#70AD47",
                     "linea":"#FF8C00","acondicionamiento":"#1D9E75"}
        labels = list(kg_mes.keys())
        values = list(kg_mes.values())
        colors = [color_map.get(k.split("_")[0], "#888") for k in labels]
        bars = ax.bar(labels, values, color=colors, edgecolor='white', linewidth=0.5)
        for bar, val in zip(bars, values):
            ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+max(values)*0.01,
                    f'{val:,.1f}kg', ha='center', fontsize=8.5, fontweight='bold')
        ax.set_ylabel('kg procesados', fontsize=9)
        plt.xticks(rotation=20, ha='right', fontsize=8)

    ax.set_title(f'RESINA TRANSFORMADA — {mes_nombre} {anio}', fontsize=12, fontweight='bold', pad=10)
    plt.tight_layout(pad=0.8)
    buf = BytesIO()
    plt.savefig(buf, format='png', dpi=120, bbox_inches='tight', facecolor='#EBEBEB')
    buf.seek(0)
    plt.close()
    return buf


def _chart_desperdicios(desp: dict, mes_nombre: str, anio: int) -> BytesIO:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt

    fig, ax = plt.subplots(figsize=(9.4, 4.8), dpi=120)
    _setup_ax(ax, fig)

    if not desp:
        ax.text(0.5, 0.5, 'Sin registros de desperdicio',
                ha='center', va='center', fontsize=12, color='#888', transform=ax.transAxes)
        ax.axis('off')
    else:
        sorted_d = sorted(desp.items(), key=lambda x: x[1], reverse=True)[:10]
        labels   = [d[0][:28] for d in sorted_d]
        values   = [d[1] for d in sorted_d]
        bars = ax.barh(labels, values, color='#4472C4', edgecolor='white')
        for bar, val in zip(bars, values):
            ax.text(val+max(values)*0.01, bar.get_y()+bar.get_height()/2,
                    f'{val:,}', va='center', fontsize=8.5, fontweight='bold')
        ax.invert_yaxis()
        ax.set_xlabel('Unidades', fontsize=9)

    ax.set_title(f'DESPERDICIO — {mes_nombre} {anio}', fontsize=12, fontweight='bold', pad=10)
    plt.tight_layout(pad=0.8)
    buf = BytesIO()
    plt.savefig(buf, format='png', dpi=120, bbox_inches='tight', facecolor='#EBEBEB')
    buf.seek(0)
    plt.close()
    return buf


def _chart_mttr(mttr_data: dict, mes_nombre: str, anio: int) -> BytesIO:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt

    fig, axes = plt.subplots(1, 3, figsize=(9.4, 4.2), dpi=120)
    fig.patch.set_facecolor('#EBEBEB')
    fig.suptitle(f'MTTR · MTBF · DISPONIBILIDAD — {mes_nombre} {anio}',
                 fontsize=10, fontweight='bold')

    label_map = {
        "inyeccion":"Inyección","soplado":"Soplado",
        "linea_copro":"L.Copro","linea_orina":"L.Orina",
        "acondicionamiento":"Acond.",
    }

    if not mttr_data:
        for ax in axes:
            ax.text(0.5, 0.5, 'Sin datos', ha='center', va='center',
                    fontsize=10, color='#888', transform=ax.transAxes)
            ax.axis('off')
    else:
        tipos  = list(mttr_data.keys())
        labels = [label_map.get(t, t) for t in tipos]
        configs = [
            (axes[0], [mttr_data[t]["mttr"] for t in tipos], '#E06C00', 'MTTR (min)', 'T. medio reparar'),
            (axes[1], [mttr_data[t]["mtbf"] for t in tipos], '#4472C4', 'MTBF (min)', 'T. medio entre fallas'),
            (axes[2], [mttr_data[t]["disp"] for t in tipos], '#70AD47', 'Disp. (%)',  'Disponibilidad'),
        ]
        for ax, vals, color, ylabel, subtitle in configs:
            _setup_ax(ax, fig)
            bars = ax.bar(range(len(labels)), vals, color=color, edgecolor='white')
            mx = max(vals) if vals else 1
            for bar, val in zip(bars, vals):
                ax.text(bar.get_x()+bar.get_width()/2,
                        bar.get_height() + mx*0.03,
                        f'{val:.1f}', ha='center', fontsize=7.5, fontweight='bold')
            ax.set_xticks(range(len(labels)))
            ax.set_xticklabels(labels, fontsize=7, rotation=20, ha='right')
            ax.set_ylabel(ylabel, fontsize=8)
            ax.set_title(subtitle, fontsize=8.5, fontweight='bold')

    plt.tight_layout(pad=1.0)
    buf = BytesIO()
    plt.savefig(buf, format='png', dpi=120, bbox_inches='tight', facecolor='#EBEBEB')
    buf.seek(0)
    plt.close()
    return buf


# ── Construcción PPTX ─────────────────────────────────────────────────────────
def _build_pptx(mes: int, anio: int, datos: dict) -> BytesIO:
    mes_nombre   = MESES_ES[mes]
    por_tipo_mes = datos["por_tipo_mes"]

    # Abrir template y eliminar todos sus slides para usar solo el master/layouts
    prs = Presentation(TEMPLATE_PATH)

    # Eliminar slides existentes manipulando XML directamente
    NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    xml_slides = prs.slides._sldIdLst
    for sldId_el in list(xml_slides):
        rid = sldId_el.get(f'{{{NS}}}id') or sldId_el.get('r:id')
        xml_slides.remove(sldId_el)
        try:
            if rid in prs.part.rels:
                prs.part.drop_rel(rid)
        except Exception:
            pass

    # Layout con fondo verde (índice 2 en este template)
    layout = prs.slide_layouts[2]

    def new_slide():
        slide = prs.slides.add_slide(layout)
        for ph in list(slide.placeholders):
            ph._element.getparent().remove(ph._element)
        return slide

    def add_logo(slide):
        if os.path.exists(LOGO_PATH):
            slide.shapes.add_picture(
                LOGO_PATH, Inches(0.25), Inches(0.10),
                Inches(2.20), Inches(0.60)
            )

    def add_title(slide, texto: str):
        tb = slide.shapes.add_textbox(Inches(0.25), Inches(0.75), Inches(12.80), Inches(0.55))
        tf = tb.text_frame
        tf.word_wrap = False
        p  = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = texto
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x21, 0x63, 0x63)
        run.font.name = "Arial"

    def add_chart(slide, buf: BytesIO, x=1.80, y=1.30, w=9.40, h=5.80):
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(x), Inches(y), Inches(w), Inches(h))

    # Slide 1 — Portada
    s = new_slide()
    add_logo(s)
    tb = s.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(10.0), Inches(2.0))
    tf = tb.text_frame
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{mes_nombre} {anio}"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x21, 0x63, 0x63)
    run.font.name = "Arial"

    # Slide 2 — OEE Inyección
    s = new_slide(); add_logo(s)
    add_title(s, f"INDICADORES INYECCION {mes_nombre} {anio}")
    add_chart(s, _chart_oee_linea("INYECCION", por_tipo_mes, "inyeccion", mes))

    # Slide 3 — OEE Soplado
    s = new_slide(); add_logo(s)
    add_title(s, f"INDICADORES SOPLADO {mes_nombre} {anio}")
    add_chart(s, _chart_oee_linea("OEE SOPLADO", por_tipo_mes, "soplado", mes))

    # Slide 4 — Línea Orina
    s = new_slide(); add_logo(s)
    add_title(s, f"INDICADORES LINEA FRASCO ORINA {mes_nombre} {anio}")
    add_chart(s, _chart_oee_linea("LINEA FRASCO ORINA", por_tipo_mes, "linea_orina", mes))

    # Slide 5 — Línea Coprológico
    s = new_slide(); add_logo(s)
    add_title(s, f"INDICADORES LINEA FRASCO COPROLOGICO {mes_nombre} {anio}")
    add_chart(s, _chart_oee_linea("LINEA FRASCO COPROLOGICO", por_tipo_mes, "linea_copro", mes))

    # Slide 6 — Acondicionamiento
    s = new_slide(); add_logo(s)
    add_title(s, f"INDICADORES ACONDICIONAMIENTO {mes_nombre} {anio}")
    add_chart(s, _chart_oee_linea("ACONDICIONAMIENTO", por_tipo_mes, "acondicionamiento", mes))

    # Slide 7 — Resina / kg
    s = new_slide(); add_logo(s)
    add_title(s, f"INDICADORES RESINA TRANSFORMADA {mes_nombre} {anio}")
    add_chart(s, _chart_kg(datos["kg_mes"], mes_nombre, anio))

    # Slide 8 — Desperdicio
    s = new_slide(); add_logo(s)
    add_title(s, f"DESPERDICIO {mes_nombre} {anio}")
    add_chart(s, _chart_desperdicios(datos["desperdicios"], mes_nombre, anio))

    # Slide 9 — MTTR/MTBF
    s = new_slide(); add_logo(s)
    add_title(s, f"MTTR, MTBF, DISPONIBILIDAD {mes_nombre} {anio}")
    add_chart(s, _chart_mttr(datos["mttr_mtbf"], mes_nombre, anio), y=1.35, h=5.70)

    # Slide 10 — Gracias
    s = new_slide(); add_logo(s)
    tb = s.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.0), Inches(2.0))
    tf = tb.text_frame
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "GRACIAS"
    run.font.size = Pt(60)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x21, 0x63, 0x63)
    run.font.name = "Arial Black"

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# ── ENDPOINTS ─────────────────────────────────────────────────────────────────

@router.get("/reporte-mensual/preview")
def preview_datos_mes(mes: int, anio: int, db: Session = Depends(get_db)):
    if not (1 <= mes <= 12):
        raise HTTPException(status_code=400, detail="Mes inválido (1-12)")

    datos = _datos_mes(mes, anio, db)
    TIPO_LABELS = {
        "inyeccion":"Inyección","soplado":"Soplado",
        "linea_copro":"Línea Copro","linea_orina":"Línea Orina",
        "acondicionamiento":"Acondicionamiento",
    }
    resultado: dict = {"mes": mes, "anio": anio, "tipos": {}}
    for tipo_key, meses_data in datos["por_tipo_mes"].items():
        mensual = []
        for m in range(1, 13):
            vals = meses_data.get(m, [])
            avg  = round(sum(vals)/len(vals), 2) if vals else None
            mensual.append({"mes": m, "oee": avg})
        resultado["tipos"][tipo_key] = {
            "label":   TIPO_LABELS.get(tipo_key, tipo_key),
            "mensual": mensual,
            "oee_mes": next((p["oee"] for p in mensual if p["mes"] == mes), None),
        }
    resultado["mttr_mtbf"]    = datos["mttr_mtbf"]
    resultado["desperdicios"] = datos["desperdicios"]
    resultado["kg_mes"]       = datos["kg_mes"]
    return resultado


@router.get("/reporte-mensual/pptx")
def generar_pptx(mes: int, anio: int, db: Session = Depends(get_db)):
    if not (1 <= mes <= 12):
        raise HTTPException(status_code=400, detail="Mes inválido (1-12)")
    if not os.path.exists(TEMPLATE_PATH):
        raise HTTPException(status_code=500,
            detail=f"Template no encontrado: {TEMPLATE_PATH}")

    datos    = _datos_mes(mes, anio, db)
    output   = _build_pptx(mes, anio, datos)
    filename = f"Reporte_{MESES_ES[mes]}_{anio}_Inverfarma.pptx"

    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'}
    )